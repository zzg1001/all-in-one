"""
OCR 服务 - 文字识别与文档内容提取
使用 DashScope Anthropic 兼容接口 (qwen3.6-plus)
"""
import io
import base64
import logging
import os

logger = logging.getLogger(__name__)

# DashScope Anthropic 兼容接口配置
DASHSCOPE_BASE_URL = os.getenv("DASHSCOPE_BASE_URL", "https://dashscope.aliyuncs.com/apps/anthropic")
DASHSCOPE_API_KEY = os.getenv("DASHSCOPE_API_KEY", "sk-f56f11bb8f0e48a985a655b36ce2970e")
DASHSCOPE_MODEL = os.getenv("DASHSCOPE_MODEL", "qwen3.6-plus")


class OCRService:
    """OCR 识别服务 - 使用 qwen3.6-plus"""

    def __init__(self):
        self.base_url = DASHSCOPE_BASE_URL
        self.api_key = DASHSCOPE_API_KEY
        self.model = DASHSCOPE_MODEL

    async def recognize(self, content: bytes, ext: str, filename: str) -> str:
        """
        识别文件内容

        Args:
            content: 文件二进制内容
            ext: 文件扩展名
            filename: 文件名

        Returns:
            识别出的文字内容
        """
        ext = ext.lower()

        # 图片 OCR
        if ext in ('jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'tiff'):
            return await self._recognize_image(content, filename)

        # PDF 文字提取
        if ext == 'pdf':
            return await self._extract_pdf(content)

        # PPT 文字提取
        if ext in ('ppt', 'pptx'):
            return await self._extract_ppt(content)

        # Word 文字提取
        if ext in ('doc', 'docx'):
            return await self._extract_word(content)

        raise ValueError(f"不支持的文件格式: {ext}")

    def _compress_image(self, content: bytes, max_size_mb: float = 4.0) -> bytes:
        """压缩图片到指定大小以下"""
        max_bytes = int(max_size_mb * 1024 * 1024)

        # 如果已经小于限制，直接返回
        if len(content) <= max_bytes:
            return content

        try:
            from PIL import Image

            img = Image.open(io.BytesIO(content))

            # 如果是 RGBA，转为 RGB
            if img.mode == 'RGBA':
                img = img.convert('RGB')

            # 尝试不同质量压缩
            for quality in [85, 70, 50, 30, 20]:
                buffer = io.BytesIO()
                img.save(buffer, format='JPEG', quality=quality, optimize=True)
                data = buffer.getvalue()
                logger.info(f"压缩图片: quality={quality}, size={len(data)} bytes")
                if len(data) <= max_bytes:
                    return data

            # 如果还是太大，缩小尺寸
            scale = 0.7
            while scale > 0.2:
                new_size = (int(img.width * scale), int(img.height * scale))
                resized = img.resize(new_size, Image.LANCZOS)
                buffer = io.BytesIO()
                resized.save(buffer, format='JPEG', quality=50)
                data = buffer.getvalue()
                logger.info(f"缩小图片: scale={scale:.2f}, size={new_size}, bytes={len(data)}")
                if len(data) <= max_bytes:
                    return data
                scale *= 0.7

            # 最后手段，返回最小的版本
            return data

        except ImportError:
            logger.warning("PIL 未安装，无法压缩图片")
            return content
        except Exception as e:
            logger.error(f"图片压缩失败: {e}")
            return content

    async def _recognize_image(self, content: bytes, filename: str = "image.jpg") -> str:
        """使用 qwen3.6-plus 进行图片 OCR 识别"""
        logger.info(f"开始 OCR 识别: {filename}, 文件大小: {len(content)} bytes")

        try:
            import anthropic

            # 压缩图片（限制 4MB）
            content = self._compress_image(content, max_size_mb=4.0)
            logger.info(f"压缩后大小: {len(content)} bytes")

            # 转为 base64
            image_data = base64.b64encode(content).decode('utf-8')

            # 压缩后统一使用 JPEG
            media_type = 'image/jpeg'

            logger.info(f"调用 DashScope qwen3.6-plus OCR, base_url={self.base_url}")

            client = anthropic.Anthropic(
                api_key=self.api_key,
                base_url=self.base_url
            )

            message = client.messages.create(
                model=self.model,
                max_tokens=8192,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": image_data
                            }
                        },
                        {
                            "type": "text",
                            "text": """请仔细识别这张图片中的所有文字内容，完整输出。

如果是营业执照，请按以下格式输出每个字段：
统一社会信用代码：xxx
名称：xxx
类型：xxx
法定代表人：xxx
注册资本：xxx
成立日期：xxx
住所：xxx
经营范围：xxx

如果是其他文档，请完整输出所有可见文字，保持原有格式。
不要遗漏任何文字内容。"""
                        }
                    ]
                }]
            )

            # 处理响应 - 提取文本内容（跳过 ThinkingBlock）
            result = ""
            if message.content:
                for block in message.content:
                    # 检查是否是 TextBlock
                    if hasattr(block, 'text'):
                        result += block.text
                    elif hasattr(block, 'type') and block.type == 'text':
                        result += getattr(block, 'text', '')

            logger.info(f"OCR 识别成功，文本长度: {len(result)} 字符")
            return result

        except Exception as e:
            logger.error(f"OCR 识别失败: {e}")
            raise ValueError(f"OCR 识别失败: {str(e)}")

    async def _extract_pdf(self, content: bytes) -> str:
        """PDF 文字提取 - 支持文字型 PDF 和扫描件"""
        try:
            # 1. 先尝试用 pdfplumber 提取文字
            extracted_text = ""
            try:
                import pdfplumber

                text_parts = []
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    for i, page in enumerate(pdf.pages, 1):
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_parts.append(f"--- 第 {i} 页 ---\n{page_text}")

                if text_parts:
                    extracted_text = '\n\n'.join(text_parts)

            except ImportError:
                logger.warning("pdfplumber 未安装")

            # 2. 如果提取的文字太少（可能是扫描件），使用 OCR
            if len(extracted_text.strip()) < 50:
                logger.info(f"PDF 提取文字过少({len(extracted_text)}字)，尝试 OCR 识别")
                ocr_text = await self._pdf_to_ocr(content)
                if ocr_text:
                    return ocr_text

            if extracted_text:
                return extracted_text

            return "PDF 中未找到文字内容"

        except Exception as e:
            logger.error(f"PDF 提取失败: {e}")
            raise ValueError(f"PDF 解析失败: {str(e)}")

    async def _pdf_to_ocr(self, content: bytes) -> str:
        """将 PDF 转为图片并调用 OCR"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(stream=content, filetype="pdf")
            all_text = []

            for page_num in range(min(len(doc), 5)):  # 最多处理前5页
                page = doc[page_num]

                # 将页面渲染为图片（分辨率 200 DPI）
                mat = fitz.Matrix(200/72, 200/72)
                pix = page.get_pixmap(matrix=mat)

                # 转为 PNG 格式
                img_bytes = pix.tobytes("png")

                logger.info(f"PDF 第 {page_num + 1} 页转图片完成，大小: {len(img_bytes)} bytes")

                # 调用 OCR 识别
                page_text = await self._recognize_image(img_bytes, f"page_{page_num + 1}.png")

                if page_text and page_text.strip():
                    all_text.append(f"--- 第 {page_num + 1} 页 ---\n{page_text}")

            doc.close()

            if all_text:
                return '\n\n'.join(all_text)

            return ""

        except ImportError:
            logger.warning("PyMuPDF (fitz) 未安装，无法对 PDF 扫描件进行 OCR")
            return ""
        except Exception as e:
            logger.error(f"PDF OCR 失败: {e}")
            return ""

    async def _extract_ppt(self, content: bytes) -> str:
        """PPT 文字提取"""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                if slide_texts:
                    text_parts.append(f"--- 幻灯片 {slide_num} ---\n" + '\n'.join(slide_texts))

            if text_parts:
                return '\n\n'.join(text_parts)

            return "PPT 中未找到文字内容"

        except ImportError:
            return "PPT 解析功能需要安装 python-pptx 库"
        except Exception as e:
            logger.error(f"PPT 提取失败: {e}")
            raise ValueError(f"PPT 解析失败: {str(e)}")

    async def _extract_word(self, content: bytes) -> str:
        """Word 文字提取"""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())

            # 也提取表格中的文字
            for table in doc.tables:
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_texts.append(cell.text.strip())
                    if row_texts:
                        text_parts.append(' | '.join(row_texts))

            if text_parts:
                return '\n\n'.join(text_parts)
            return "Word 文档中未找到文字内容"

        except ImportError:
            return "Word 解析功能需要安装 python-docx 库"
        except Exception as e:
            logger.error(f"Word 提取失败: {e}")
            raise ValueError(f"Word 解析失败: {str(e)}")
